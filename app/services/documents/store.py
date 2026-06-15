"""
Locate, read and persist generated documents — scoped per project.

Every deliverable is written under a per-project folder
(``GENERATED_DIR/project_<id>``) so a document only ever shows up in the project
that created it. Alongside each file we keep:

  * a *source sidecar* (the original Markdown or Python the document was built
    from) under a hidden ``.sources`` subfolder, so a document can be re-read and
    modified later without lossy text extraction from the binary, and
  * a *metadata sidecar* (JSON: title, creation timestamp, format, project) under
    a hidden ``.meta`` subfolder, surfaced in the UI's files panel.

Both subfolders are hidden so they never show up in the downloadable-files list.
"""
from __future__ import annotations

import json
import re
from pathlib import Path
from datetime import datetime, timezone

from app.config import GENERATED_DIR
from app.logging_config import logger


def project_dir(project_id: int) -> Path:
    """Return (and create) the per-project folder that holds its deliverables."""
    d = GENERATED_DIR / f"project_{int(project_id)}"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _sources_dir(project_id: int) -> Path:
    d = project_dir(project_id) / ".sources"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _meta_dir(project_id: int) -> Path:
    d = project_dir(project_id) / ".meta"
    d.mkdir(parents=True, exist_ok=True)
    return d


def build_dest(project_id: int, formato: str, nombre: str) -> Path:
    """Return a unique, sanitized destination path inside the project's folder."""
    fmt = formato.lower().strip()
    safe_name = re.sub(r"[^\w\-]", "_", nombre) or "document"
    timestamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    return project_dir(project_id) / f"{safe_name}_{timestamp}.{fmt}"


def save_source_sidecar(project_id: int, filename: str, source: str) -> None:
    """Persist the source (Markdown/Python) a document was generated from."""
    try:
        (_sources_dir(project_id) / f"{filename}.source.md").write_text(
            source, encoding="utf-8"
        )
    except Exception:  # pragma: no cover - sidecar is best-effort
        logger.warning("Could not save the source sidecar for %s", filename)


def save_metadata(
    project_id: int,
    filename: str,
    *,
    title: str,
    formato: str,
) -> dict:
    """Persist and return metadata for a generated document.

    Stored fields: title, created_at (UTC ISO-8601), format, project_id, filename.
    """
    meta = {
        "filename": filename,
        "title": (title or "").strip() or filename,
        "format": formato.lower().strip(),
        "project_id": int(project_id),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    try:
        (_meta_dir(project_id) / f"{filename}.meta.json").write_text(
            json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8"
        )
    except Exception:  # pragma: no cover - sidecar is best-effort
        logger.warning("Could not save the metadata sidecar for %s", filename)
    return meta


def read_metadata(project_id: int, filename: str) -> dict | None:
    """Return the stored metadata for a generated document, or None."""
    path = _meta_dir(project_id) / f"{filename}.meta.json"
    if not path.is_file():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:  # pragma: no cover
        return None


def list_generated(project_id: int) -> list[dict]:
    """List a project's deliverables (newest first) with size + metadata.

    Each entry: {filename, size, title, created_at, format}. Hidden sidecar
    folders are skipped. Files with no metadata sidecar fall back to filesystem
    info so legacy files still appear.
    """
    pdir = project_dir(project_id)
    if not pdir.exists():
        return []
    entries: list[dict] = []
    for f in pdir.iterdir():
        if not f.is_file() or f.name.startswith("."):
            continue
        meta = read_metadata(project_id, f.name) or {}
        stat = f.stat()
        entries.append(
            {
                "filename": f.name,
                "size": stat.st_size,
                "title": meta.get("title", f.name),
                "created_at": meta.get(
                    "created_at",
                    datetime.fromtimestamp(stat.st_mtime, timezone.utc).isoformat(),
                ),
                "format": meta.get("format", f.suffix.lstrip(".").lower()),
                "mtime": stat.st_mtime,
            }
        )
    entries.sort(key=lambda e: e.pop("mtime"), reverse=True)
    return entries


def safe_generated_path(project_id: int, filename: str) -> Path | None:
    """Resolve a download/delete target inside the project's folder, or None."""
    if "/" in filename or "\\" in filename or ".." in filename:
        return None
    path = project_dir(project_id) / filename
    return path if path.is_file() else None


def delete_generated(project_id: int, filename: str) -> bool:
    """Delete a deliverable and its sidecars. Returns True if the file existed."""
    path = safe_generated_path(project_id, filename)
    if path is None:
        return False
    path.unlink()
    for side in (
        _sources_dir(project_id) / f"{filename}.source.md",
        _meta_dir(project_id) / f"{filename}.meta.json",
    ):
        try:
            side.unlink(missing_ok=True)
        except Exception:  # pragma: no cover
            pass
    return True


def _resolve_generated(project_id: int, nombre: str) -> Path | None:
    """Locate a generated file by name (with or without extension) in the project.

    If several files match (same name, different timestamp) the most recent one
    is returned.
    """
    pdir = project_dir(project_id)
    base = re.sub(r"\.source\.md$", "", nombre.strip())
    base = Path(base).name  # never accept paths
    candidates: list[Path] = []
    if (pdir / base).is_file():
        candidates.append(pdir / base)
    stem = re.sub(r"\.[^.]+$", "", base)
    for p in pdir.glob(f"{stem}*"):
        if p.is_file() and not p.name.endswith(".source.md"):
            candidates.append(p)
    if not candidates:
        return None
    return max(set(candidates), key=lambda p: p.stat().st_mtime)


def read_generated(project_id: int, nombre: str) -> str:
    """Return the text content of an already-generated document in the project.

    Prefers the source sidecar (faithful re-read of what produced the file); if
    absent, extracts text from the binary according to its format.
    """
    target = _resolve_generated(project_id, nombre)
    if target is None:
        return (
            f"No generated document named '{nombre}' was found. "
            "Check the exact file name."
        )

    sidecar = _sources_dir(project_id) / f"{target.name}.source.md"
    if sidecar.is_file():
        return sidecar.read_text(encoding="utf-8")

    fmt = target.suffix.lower().lstrip(".")
    try:
        if fmt in ("md", "txt", "html"):
            return target.read_text(encoding="utf-8")

        if fmt == "pdf":
            import fitz  # PyMuPDF

            doc = fitz.open(str(target))
            return "\n".join(page.get_text() for page in doc)

        if fmt == "docx":
            from docx import Document

            doc = Document(str(target))
            return "\n".join(p.text for p in doc.paragraphs)

        if fmt == "xlsx":
            from openpyxl import load_workbook

            wb = load_workbook(str(target), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"# {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    lines.append("\t".join("" if c is None else str(c) for c in row))
            return "\n".join(lines)

        if fmt == "pptx":
            from pptx import Presentation

            prs = Presentation(str(target))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        lines.append(shape.text_frame.text)
            return "\n".join(lines)

    except Exception as exc:  # pragma: no cover - extraction is best-effort
        logger.exception("Error reading generated document '%s': %s", target.name, exc)
        return (
            f"Could not extract the text from '{target.name}' (format {fmt}): {exc}. "
            "If you need to modify it, regenerate it from the source content."
        )

    return target.read_text(encoding="utf-8", errors="replace")
